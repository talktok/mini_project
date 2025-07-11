from pptx import Presentation
from dotenv import load_dotenv
from openai import OpenAI
import os
import json
import tiktoken

# ✅ 환경변수 불러오기
load_dotenv(override=True)
apikey = os.getenv('OPEN_API_KEY')
model = os.getenv("GET_DEFAULT_MODEL", "gpt-4o-mini")

# ✅ 토큰 인코딩 초기화
encoding = tiktoken.encoding_for_model(model)

# ✅ OpenAI 클라이언트 설정
client = OpenAI(api_key=apikey)

# ✅ PPTX 슬라이드 텍스트 추출 함수
def extract_slide_texts(pptx_path):
    prs = Presentation(pptx_path)
    texts = []
    for slide in prs.slides:
        text = ""
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text.strip() + "\n"
        texts.append(text.strip())
    return texts

# ✅ 슬라이드별 키워드 추출 함수
def extract_keywords_per_slide(slide_texts):
    keywords_dict = {}
    for i, text in enumerate(slide_texts, 1):
        prompt = """
        당신은 정보를 핵심 포인트로 전달하는 데 특화된 능숙한 AI입니다.
        다음 텍스트를 기반으로 논의되거나 언급된 주요 포인트를 확인하고 문장이 아닌 짧은 단어 또는 어절의 형태로 최대 5개를 나열합니다. 
        이는 논의의 본질에 가장 중요한 아이디어, 결과 또는 주제가 되어야 합니다.
        당신의 목표는 누군가가 읽을 수 있는 목록을 제공하여 이야기된 내용을 빠르게 이해하는 것입니다.
        """
        response = client.chat.completions.create(
            model=model,
            messages=[
                {"role": "system", "content": prompt},
                {"role": "user", "content": text}
            ],
            temperature=0
        )
        raw_lines = response.choices[0].message.content.strip().splitlines()
        keywords = [line.lstrip("-•●").strip() for line in raw_lines if line.strip()]
        keywords_dict[i] = keywords
    return keywords_dict

# ✅ 슬라이드 흐름 기반 발표 대본 생성
def presentation_scripts(slides_dict, system_prompt):
    scripts = {}
    assistant_prompt = ''

    for slide_num, slide_text in slides_dict.items():
        user_prompt = f"[슬라이드 내용]: {slide_text}"
        messages = [{"role": "system", "content": system_prompt}]
        if assistant_prompt:
            messages.append({"role": "assistant", "content": assistant_prompt})
        messages.append({"role": "user", "content": user_prompt})

        response = client.chat.completions.create(
            model=model,
            messages=messages,
            temperature=0,
            max_tokens=500
        )
        script = response.choices[0].message.content.strip()
        assistant_prompt += script + "\n"
        scripts[slide_num] = script

        tokens = encoding.encode(script)
        print(f"📦 슬라이드 {slide_num} 토큰 수: {len(tokens)}")
    return scripts

# ✅ 전체 대본을 매끄럽게 다시 쓰는 함수
def polish_final_script(presenter, audience, purpose, tone, time, raw_script):
    polish_prompt = f"""
    당신은 발표에 굉장히 능숙하고 숙련된 {presenter} 입니다.
    당신의 발표를 듣는 발표대상은 {audience} 입니다.
    발표목적은 {purpose} 이며, 말투는 {tone} 말합니다.
    발표시간의 시간은 {time}분 이므로, 시간에 맞게 분량을 설정합니다.
    다음 대본을 지정된 설정을 따르면서 발표 대본을 생성합니다.
    흐름이 자연스러우며 핵심 메시지를 분명하게 전달하고, 복잡한 내용을 간결하게 설명합니다.
    당신은 자신감과 진정성 있는 태도를 갖추고 청중과 소통하며 흥미를 유발해 집중을 잘 끌어내고 핵심 내용을 강조하여 전달력이 명확해야 합니다.
    """

    response = client.chat.completions.create(
        model=model,
        temperature=0.7,
        messages=[
            {"role": "system", "content": polish_prompt},
            {"role": "user", "content": raw_script}
        ]
    )
    return response.choices[0].message.content.strip()

# ✅ 메인 실행 함수
def main():
    pptx_path = "ex.pptx"  # ⚠️ 실제 PPTX 경로로 수정해줘

    # 1. 슬라이드 텍스트 추출
    slide_texts = extract_slide_texts(pptx_path)

    # 2. 키워드 추출
    keywords_dict = extract_keywords_per_slide(slide_texts)

    # 3. 딕셔너리로 변환
    slides_dict = {i + 1: text for i, text in enumerate(slide_texts)}

    # 4. 시스템 프롬프트 정의
    system_prompt = (
        """
        넌 발표에 능숙하고 숙련된 전문가 AI야. 각 슬라이드 내용을 기반으로 자연스럽고 발표 톤으로 2~3문장 대본을 작성해줘.
        이전 슬라이드의 흐름을 고려해서 자연스럽게 이어지도록 해줘.
        핵심 내용을 강조하여 전달력이 명확해야 합니다.
        분량은 각 슬라이드당 한 문단으로 만들어야 합니다.
        """
    )

    print("\n🚀 슬라이드별 키워드 및 대본 생성 중...\n")
    generated_scripts = presentation_scripts(slides_dict, system_prompt)

    # ✅ 슬라이드별 키워드 + 대본 출력
    for slide_num in sorted(slides_dict):
        print("=" * 50)
        print(f"🟨 [슬라이드 {slide_num}]")
        print(f"📜 텍스트:\n{slides_dict[slide_num]}")
        print(f"🔑 키워드: {', '.join(keywords_dict[slide_num])}")
        print(f"🎙️ 대본:\n{generated_scripts[slide_num]}\n")

    # ✅ 전체 발표 대본 (초안)
    raw_script = "\n".join(generated_scripts[slide_num] for slide_num in sorted(generated_scripts))

    # ✅ OpenAI로 자연스럽게 정리된 전체 발표 대본 출력
    polished_script = polish_final_script(presenter, audience, purpose, tone, time, raw_script)
    print("\n✅ 최종 발표 대본\n" + "=" * 30)
    print(polished_script)

if __name__ == "__main__":
    main()
