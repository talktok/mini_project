from flask import Flask, render_template, request, redirect, url_for, flash, session, send_file, jsonify
from flask_sqlalchemy import SQLAlchemy
from flask_login import LoginManager, UserMixin, login_user, logout_user, login_required, current_user
from werkzeug.security import generate_password_hash, check_password_hash
from werkzeug.utils import secure_filename
import os
import uuid
import random # For dummy AI evaluation data

# PPTX 및 OpenAI 관련 라이브러리 임포트
from pptx import Presentation
from dotenv import load_dotenv
from openai import OpenAI
import json
import tiktoken

# ✅ 음성 인식 및 텍스트 유사도 관련 라이브러리 임포트
import torch
from transformers import AutoModelForSpeechSeq2Seq, AutoProcessor, pipeline
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity
from pydub import AudioSegment, silence # pydub 추가 임포트

# ✅ 환경변수 불러오기 (OpenAI API 키)
load_dotenv(override=True)
OPENAI_API_KEY = os.getenv('OPEN_API_KEY')
DEFAULT_MODEL = os.getenv("GET_DEFAULT_MODEL", "gpt-4o-mini") # 기본 모델 설정

# ✅ OpenAI 클라이언트 설정
client = OpenAI(api_key=OPENAI_API_KEY)

# ✅ 토큰 인코딩 초기화
encoding = tiktoken.encoding_for_model(DEFAULT_MODEL)

app = Flask(__name__)

app.config['SECRET_KEY'] = 'your_super_secret_key_for_session_security'
app.config['SQLALCHEMY_DATABASE_URI'] = 'sqlite:///site.db'
app.config['SQLALCHEMY_TRACK_MODIFICATIONS'] = False

db = SQLAlchemy(app)
login_manager = LoginManager(app)
login_manager.login_view = 'login'

UPLOAD_FOLDER = 'uploads'
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER

# 허용된 파일 확장자
ALLOWED_EXTENSIONS = {'ppt', 'pptx'}
ALLOWED_AUDIO_EXTENSIONS = {'mp3', 'wav', 'webm', 'm4a'} # m4a 추가

# 최대 파일 크기 (PPT: 16MB, Audio: 50MB)
app.config['MAX_CONTENT_LENGTH'] = 50 * 1024 * 1024  # Flask 전체 요청 크기 제한 (오디오 파일에 맞춰 설정)
app.config['MAX_PPT_FILE_SIZE'] = 16 * 1024 * 1024   # PPT 파일 개별 크기 제한
app.config['MAX_AUDIO_FILE_SIZE'] = 50 * 1024 * 1024 # 오디오 파일 개별 크기 제한

# ✅ Whisper 모델 준비 (앱 시작 시 한 번만 로드)
# GPU 사용 가능 여부 확인
device = "cuda:0" if torch.cuda.is_available() else "cpu"
model_id = "openai/whisper-small" # 더 작은 모델로 변경하여 로딩 시간 단축 및 메모리 절약
# model_id = "openai/whisper-base" # 필요시 더 큰 모델 사용 가능
# model_id = "openai/whisper-tiny" # 가장 작은 모델

print(f"Loading Whisper model on device: {device}")
try:
    processor = AutoProcessor.from_pretrained(model_id)
    model = AutoModelForSpeechSeq2Seq.from_pretrained(model_id).to(device)
    pipe = pipeline(
        "automatic-speech-recognition",
        model=model,
        tokenizer=processor.tokenizer,
        feature_extractor=processor.feature_extractor,
        device=0 if device.startswith("cuda") else -1, # device 인덱스 지정 (cuda:0 -> 0)
        return_timestamps=False # 타임스탬프는 필요 없으므로 False
    )
    print("Whisper model loaded successfully.")
except Exception as e:
    print(f"ERROR: Failed to load Whisper model: {e}")
    pipe = None # 모델 로드 실패 시 pipe를 None으로 설정하여 이후 오류 방지

class User(db.Model, UserMixin):
    id = db.Column(db.Integer, primary_key=True)
    username = db.Column(db.String(20), unique=True, nullable=False)
    email = db.Column(db.String(120), unique=True, nullable=False)
    password_hash = db.Column(db.String(128), nullable=False)

    def set_password(self, password):
        self.password_hash = generate_password_hash(password)

    def check_password(self, password):
        return check_password_hash(self.password_hash, password)

    def get_id(self):
        return str(self.id)

    def __repr__(self):
        return f"User('{self.username}', '{self.email}')"

@login_manager.user_loader
def load_user(user_id):
    return db.session.get(User, int(user_id))

def allowed_file(filename):
    """허용된 파일 확장자인지 확인"""
    return '.' in filename and \
           filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def allowed_audio_file(filename):
    """허용된 오디오 파일 확장자인지 확인"""
    return '.' in filename and \
           filename.rsplit('.', 1)[1].lower() in ALLOWED_AUDIO_EXTENSIONS

# ✅ 오디오 파일을 WAV로 변환하는 함수
def convert_to_wav(input_path, output_path):
    try:
        audio = AudioSegment.from_file(input_path)
        audio.export(output_path, format="wav")
        return True
    except Exception as e:
        print(f"Error converting audio to WAV: {e}")
        return False

# ✅ 텍스트 유사도 계산 함수 (코사인 유사도)
def calc_similarity(text1, text2):
    # 빈 문자열 처리
    if not text1.strip() or not text2.strip():
        return 0.0 # 둘 중 하나라도 빈 문자열이면 유사도 0

    vectorizer = TfidfVectorizer().fit([text1, text2])
    vectors = vectorizer.transform([text1, text2])
    sim_score = cosine_similarity(vectors[0], vectors[1])[0][0]
    return round(sim_score * 100, 2)

# --- PPTX 슬라이드 텍스트 추출 함수 ---
def _extract_slide_texts(pptx_path):
    prs = Presentation(pptx_path)
    texts = []
    for slide in prs.slides:
        text = ""
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text.strip() + "\n"
        texts.append(text.strip())
    return texts

# --- 슬라이드별 키워드 추출 함수 (LLM 사용) ---
def _extract_keywords_per_slide(slide_texts):
    keywords_dict = {}
    for i, text in enumerate(slide_texts, 1):
        prompt = """
        당신은 정보를 핵심 포인트로 전달하는 데 특화된 능숙한 AI입니다.
        다음 텍스트를 기반으로 논의되거나 언급된 주요 포인트를 확인하고 문장이 아닌 짧은 단어 또는 어절의 형태로 최대 5개를 나열합니다. 
        이는 논의의 본질에 가장 중요한 아이디어, 결과 또는 주제가 되어야 합니다.
        당신의 목표는 누군가가 읽을 수 있는 목록을 제공하여 이야기된 내용을 빠르게 이해하는 것입니다.
        """
        try:
            response = client.chat.completions.create(
                model=DEFAULT_MODEL,
                messages=[
                    {"role": "system", "content": prompt},
                    {"role": "user", "content": text}
                ],
                temperature=0
            )
            raw_lines = response.choices[0].message.content.strip().splitlines()
            keywords = [line.lstrip("-•●").strip() for line in raw_lines if line.strip()]
            keywords_dict[i] = keywords
        except Exception as e:
            print(f"ERROR: Keyword extraction failed for slide {i}: {e}")
            keywords_dict[i] = ["키워드 추출 실패"] # 오류 시 대체 키워드
    return keywords_dict

# --- 슬라이드 흐름 기반 발표 대본 생성 (LLM 사용) ---
def _presentation_scripts(slides_dict, system_prompt):
    scripts = {}
    assistant_prompt = ''

    for slide_num, slide_text in slides_dict.items():
        user_prompt = f"[슬라이드 내용]: {slide_text}"
        messages = [{"role": "system", "content": system_prompt}]
        if assistant_prompt:
            messages.append({"role": "assistant", "content": assistant_prompt})
        messages.append({"role": "user", "content": user_prompt})

        try:
            response = client.chat.completions.create(
                model=DEFAULT_MODEL,
                messages=messages,
                temperature=0,
                max_tokens=500
            )
            script = response.choices[0].message.content.strip()
            assistant_prompt += script + "\n"
            scripts[slide_num] = script

            tokens = encoding.encode(script)
            print(f"📦 슬라이드 {slide_num} 토큰 수: {len(tokens)}")
        except Exception as e:
            print(f"ERROR: Script generation failed for slide {slide_num}: {e}")
            scripts[slide_num] = "대본 생성 중 오류가 발생했습니다." # 오류 시 대체 대본
    return scripts

# --- 전체 대본을 매끄럽게 다시 쓰는 함수 (LLM 사용) ---
def _polish_final_script(user_profile, raw_script):
    presenter = user_profile.get('presenter', '발표자')
    audience = user_profile.get('audience', '청중')
    purpose = user_profile.get('purpose', '정보 전달')
    tone = user_profile.get('tone', '부드럽고 자연스럽게')
    time = user_profile.get('time', '5') # 분 단위로 가정

    polish_prompt = f"""
    당신은 발표에 굉장히 능숙하고 숙련된 {presenter} 입니다.
    당신의 발표를 듣는 발표대상은 {audience} 입니다.
    발표목적은 {purpose} 이며, 말투는 {tone} 말합니다.
    발표시간의 시간은 {time}분 이므로, 시간에 맞게 분량을 설정합니다.
    다음 대본을 지정된 설정을 따르면서 발표 대본을 생성합니다.
    흐름이 자연스러우며 핵심 메시지를 분명하게 전달하고, 복잡한 내용을 간결하게 설명합니다.
    당신은 자신감과 진정성 있는 태도를 갖추고 청중과 소통하며 흥미를 유발해 집중을 잘 끌어내고 핵심 내용을 강조하여 전달력이 명확해야 합니다.
    """
    try:
        response = client.chat.completions.create(
            model=DEFAULT_MODEL,
            temperature=0.7,
            messages=[
                {"role": "system", "content": polish_prompt},
                {"role": "user", "content": raw_script}
            ]
        )
        return response.choices[0].message.content.strip()
    except Exception as e:
        print(f"ERROR: Final script polishing failed: {e}")
        return "최종 대본을 다듬는 중 오류가 발생했습니다."

# --- 메인 페이지 라우트 ---
@app.route('/')
def index():
    return render_template('index.html', is_authenticated=current_user.is_authenticated)

# --- 회원가입 라우트 ---
@app.route('/register', methods=['GET', 'POST'])
def register():
    if current_user.is_authenticated:
        return redirect(url_for('dashboard'))

    if request.method == 'POST':
        username = request.form['username']
        email = request.form['email']
        password = request.form['password']

        existing_user_username = db.session.execute(db.select(User).filter_by(username=username)).scalar_one_or_none()
        existing_user_email = db.session.execute(db.select(User).filter_by(email=email)).scalar_one_or_none()

        if existing_user_username:
            flash('사용자명이 이미 존재합니다. 다른 사용자명을 선택해주세요.', 'danger')
            return redirect(url_for('register'))
        if existing_user_email:
            flash('이메일이 이미 사용 중입니다. 다른 이메일을 사용해주세요.', 'danger')
            return redirect(url_for('register'))

        new_user = User(username=username, email=email)
        new_user.set_password(password)
        db.session.add(new_user)
        db.session.commit()
        flash('회원가입이 완료되었습니다! 이제 로그인할 수 있습니다.', 'success')
        return redirect(url_for('login'))
    return render_template('register.html')

# --- 로그인 라우트 ---
@app.route('/login', methods=['GET', 'POST'])
def login():
    if current_user.is_authenticated:
        return redirect(url_for('dashboard'))

    if request.method == 'POST':
        username = request.form['username']
        password = request.form['password']
        remember = True if request.form.get('remember_me') else False

        user = db.session.execute(db.select(User).filter_by(username=username)).scalar_one_or_none()

        if user and user.check_password(password):
            login_user(user, remember=remember)
            flash('로그인되었습니다!', 'success')
            next_page = request.args.get('next')
            return redirect(next_page or url_for('dashboard'))
        else:
            flash('로그인 실패. 사용자명 또는 비밀번호를 확인해주세요.', 'danger')
    return render_template('login.html')

# --- 로그아웃 라우트 ---
@app.route('/logout')
@login_required
def logout():
    logout_user()
    flash('로그아웃되었습니다.', 'info')
    return redirect(url_for('index'))

# --- 대시보드 라우트 (로그인 필요) ---
@app.route('/dashboard')
@login_required
def dashboard():
    return render_template('dashboard.html')

# --- 스크립트 생성 흐름 시작 (사용자 정보 입력) ---
@app.route('/script_generator/user_info', methods=['GET', 'POST'])
def user_info_form():
    is_fresh_start = False

    if request.args.get('reset'):
        # 모든 세션 데이터 초기화
        session.pop('user_profile', None)
        session.pop('ppt_filename_for_script', None)
        session.pop('ppt_original_filename', None)
        session.pop('slide_texts', None)
        session.pop('llm_suggested_keywords_by_slide', None)
        session.pop('current_slide_index', None)
        session.pop('all_selected_keywords', None)
        session.pop('script_options', None)
        session.pop('generated_script_text', None)
        session.pop('recorded_audio_filename', None)
        session.pop('recorded_audio_original_filename', None)
        is_fresh_start = True

    if request.method == 'POST':
        presenter = request.form.get('presenter')
        audience = request.form.get('audience')
        purpose = request.form.get('purpose')
        tone = request.form.get('tone')
        time = request.form.get('time')

        if not all([presenter, audience, purpose, tone, time]):
            flash('모든 사용자 정보를 입력해주세요.', 'danger')
            return render_template('user_info_form.html', user_profile=request.form, is_fresh_start=is_fresh_start)

        session['user_profile'] = {
            'presenter': presenter,
            'audience': audience,
            'purpose': purpose,
            'tone': tone,
            'time': time
        }
        flash('사용자 정보가 저장되었습니다. 이제 발표 자료 (PPT)를 업로드해주세요.', 'success')
        return redirect(url_for('upload_ppt_for_script'))

    user_profile = session.get('user_profile', {})
    return render_template('user_info_form.html', user_profile=user_profile, is_fresh_start=is_fresh_start)

# --- 스크립트 생성을 위한 PPT 업로드 (향상된 버전 및 LLM 연동) ---
@app.route('/script_generator/upload_ppt_for_script', methods=['GET', 'POST'])
def upload_ppt_for_script():
    if 'user_profile' not in session:
        flash('먼저 사용자 정보를 입력해주세요.', 'danger')
        return redirect(url_for('user_info_form'))

    if request.method == 'POST':
        if 'file' not in request.files:
            flash('파일이 업로드되지 않았습니다.', 'danger')
            return redirect(url_for('upload_ppt_for_script'))

        file = request.files['file']
        
        if file.filename == '':
            flash('파일을 선택해주세요.', 'danger')
            return redirect(url_for('upload_ppt_for_script'))

        if not allowed_file(file.filename):
            flash('PPT 파일 (.ppt, .pptx)만 업로드할 수 있습니다.', 'danger')
            return redirect(url_for('upload_ppt_for_script'))

        # 파일 크기 검증
        if file.content_length and file.content_length > app.config['MAX_PPT_FILE_SIZE']:
            flash(f'PPT 파일 크기가 너무 큽니다. 최대 {app.config["MAX_PPT_FILE_SIZE"] / (1024*1024):.0f}MB까지 업로드 가능합니다.', 'danger')
            return redirect(url_for('upload_ppt_for_script'))

        filepath = None # 파일 경로 초기화
        try:
            original_filename = secure_filename(file.filename)
            file_extension = os.path.splitext(original_filename)[1]
            unique_filename = str(uuid.uuid4()) + file_extension
            filepath = os.path.join(app.config['UPLOAD_FOLDER'], unique_filename)
            
            file.save(filepath)
            
            if not os.path.exists(filepath):
                flash('파일 저장 중 오류가 발생했습니다. 다시 시도해주세요.', 'danger')
                return redirect(url_for('upload_ppt_for_script'))
            
            # ✅ PPT 텍스트 추출 및 키워드 생성
            slide_texts = _extract_slide_texts(filepath)
            llm_suggested_keywords_by_slide = _extract_keywords_per_slide(slide_texts)

            # 세션에 저장 (JSON 직렬화 가능한 형태로)
            session['ppt_filename_for_script'] = unique_filename
            session['ppt_original_filename'] = original_filename
            session['slide_texts'] = slide_texts # 리스트는 JSON 직렬화 가능
            # 딕셔너리의 키가 숫자인 경우 JSON 직렬화 시 문자열로 변환될 수 있으므로 주의
            session['llm_suggested_keywords_by_slide'] = {str(k): v for k, v in llm_suggested_keywords_by_slide.items()}
            session['current_slide_index'] = 0 # 키워드 선택 시작 인덱스 초기화
            session['all_selected_keywords'] = [] # 사용자가 선택할 키워드 리스트 초기화

            flash(f'"{original_filename}" 파일이 성공적으로 업로드되었습니다. 이제 슬라이드별 키워드를 선택해주세요.', 'success')
            return redirect(url_for('keyword_select'))
            
        except Exception as e:
            flash(f'파일 처리 중 오류가 발생했습니다: {e}', 'danger')
            print(f"Error during PPT processing: {e}") # 서버 콘솔에 상세 오류 출력
            # 오류 발생 시 업로드된 파일 삭제 (선택 사항)
            if filepath and os.path.exists(filepath):
                os.remove(filepath)
            return redirect(url_for('upload_ppt_for_script'))

    return render_template('script_generator.html')

# --- 파일 크기 초과 오류 처리 ---
@app.errorhandler(413)
def too_large(e):
    flash(f'파일 크기가 너무 큽니다. PPT 파일은 최대 {app.config["MAX_PPT_FILE_SIZE"] / (1024*1024):.0f}MB, 오디오 파일은 최대 {app.config["MAX_AUDIO_FILE_SIZE"] / (1024*1024):.0f}MB까지 업로드 가능합니다.', 'danger')
    return redirect(request.referrer or url_for('index'))

# --- 스크립트 생성 (슬라이드별 키워드 선택) ---
@app.route('/script_generator/keywords', methods=['GET', 'POST'])
def keyword_select():
    # 필요한 세션 데이터 확인
    if 'user_profile' not in session or \
       'ppt_filename_for_script' not in session or \
       'slide_texts' not in session or \
       'llm_suggested_keywords_by_slide' not in session:
        flash('스크립트 생성에 필요한 정보가 부족합니다. 처음부터 다시 시작해주세요.', 'danger')
        return redirect(url_for('user_info_form', reset=True))

    slide_texts = session.get('slide_texts')
    llm_suggested_keywords_by_slide_str_keys = session.get('llm_suggested_keywords_by_slide')
    # 키가 문자열로 저장되었으므로 다시 숫자로 변환
    llm_suggested_keywords_by_slide = {int(k): v for k, v in llm_suggested_keywords_by_slide_str_keys.items()}

    TOTAL_SLIDES = len(slide_texts) # 동적으로 총 슬라이드 수 설정

    if 'current_slide_index' not in session or request.args.get('reset'):
        session['current_slide_index'] = 0
        session['all_selected_keywords'] = []
        flash('슬라이드별 키워드를 선택해주세요.', 'info')

    current_slide_index = session.get('current_slide_index')
    
    if current_slide_index >= TOTAL_SLIDES:
        flash('모든 슬라이드의 키워드 선택이 완료되었습니다.', 'success')
        return redirect(url_for('script_option'))

    # 현재 슬라이드에 대한 LLM 제안 키워드를 가져옴
    current_slide_keywords = llm_suggested_keywords_by_slide.get(current_slide_index + 1, [])
    if not current_slide_keywords: # LLM이 키워드를 생성하지 못했을 경우 대비
        current_slide_keywords = ["내용 없음", "키워드 없음"]


    if request.method == 'POST':
        selected_keywords_for_current_slide = request.form.getlist('keywords')
        if not selected_keywords_for_current_slide:
            flash('현재 슬라이드에 대한 키워드를 하나 이상 선택해주세요.', 'danger')
            return render_template('keyword_select.html',
                                   current_slide_index=current_slide_index,
                                   total_slides=TOTAL_SLIDES,
                                   keywords=current_slide_keywords)
        
        all_selected_keywords = session.get('all_selected_keywords', [])
        all_selected_keywords.append(selected_keywords_for_current_slide)
        session['all_selected_keywords'] = all_selected_keywords
        
        session['current_slide_index'] += 1
        flash(f'슬라이드 {current_slide_index + 1}의 키워드가 저장되었습니다.', 'success')
        return redirect(url_for('keyword_select'))
    
    return render_template('keyword_select.html',
                           current_slide_index=current_slide_index,
                           total_slides=TOTAL_SLIDES,
                           keywords=current_slide_keywords)

# --- 스크립트 옵션 선택 (서론/본론/결론) ---
@app.route('/script_generator/options', methods=['GET', 'POST'])
def script_option():
    # 필요한 세션 데이터 확인
    if 'user_profile' not in session or \
       'all_selected_keywords' not in session or \
       'slide_texts' not in session: # slide_texts도 필요
        flash('스크립트 생성에 필요한 정보가 부족합니다. 처음부터 다시 시작해주세요.', 'danger')
        return redirect(url_for('user_info_form', reset=True))

    slide_texts = session.get('slide_texts')
    # 모든 슬라이드의 키워드 선택이 완료되었는지 확인 (선택된 키워드 리스트의 길이가 총 슬라이드 수와 같아야 함)
    if len(session['all_selected_keywords']) < len(slide_texts):
        flash('모든 슬라이드의 키워드 선택을 완료해주세요.', 'danger')
        return redirect(url_for('keyword_select'))


    all_selected_keywords_by_slide = session.get('all_selected_keywords')
    # all_selected_keywords_by_slide는 이미 슬라이드별 리스트의 리스트이므로 flat하게 만들 필요 없음
    # 하지만 템플릿에 전달할 때는 flat한 리스트도 유용할 수 있으니 유지
    selected_keywords_flat = [item for sublist in all_selected_keywords_by_slide for item in sublist]
    
    intro_options = ['흥미로운 질문으로 시작', '최신 트렌드 언급', '문제 제기']
    body_options = ['사례 연구 포함', '기술적 설명 강조', '미래 전망 제시']
    conclusion_options = ['핵심 요약 및 제언', '청중에게 질문 던지기', '긍정적 비전 제시']

    if request.method == 'POST':
        intro_option = request.form.get('intro_option')
        body_option = request.form.get('body_option')
        conclusion_option = request.form.get('conclusion_option')

        if not all([intro_option, body_option, conclusion_option]):
            flash('서론, 본론, 결론 옵션을 모두 선택해주세요.', 'danger')
            return render_template('script_option.html',
                                   selected_keywords=selected_keywords_flat, # 템플릿에 필요
                                   intro_options=intro_options,
                                   body_options=body_options,
                                   conclusion_options=conclusion_options)

        session['script_options'] = {
            'intro': intro_option,
            'body': body_option,
            'conclusion': conclusion_option
        }
        flash('스크립트 옵션이 저장되었습니다. 이제 대본을 생성합니다.', 'success')
        return redirect(url_for('script_result'))

    return render_template('script_option.html', 
                           selected_keywords=selected_keywords_flat,
                           intro_options=intro_options,
                           body_options=body_options,
                           conclusion_options=conclusion_options)

# --- 완성된 대본 보여주기 (LLM 연동) ---
@app.route('/script_generator/result')
def script_result():
    # 세션 지속성 설정
    session.permanent = True
    
    user_profile = session.get('user_profile')
    all_selected_keywords_by_slide = session.get('all_selected_keywords')
    script_options = session.get('script_options')
    slide_texts = session.get('slide_texts')  # 원본 슬라이드 텍스트 필요

    if not user_profile or not all_selected_keywords_by_slide or not script_options or not slide_texts:
        flash('스크립트 생성에 필요한 정보가 부족합니다. 처음부터 다시 시작해주세요.', 'danger')
        return redirect(url_for('user_info_form', reset=True))

    # 슬라이드 텍스트를 딕셔너리 형태로 변환 (LLM 함수에 전달하기 위함)
    slides_dict = {i + 1: text for i, text in enumerate(slide_texts)}

    # 시스템 프롬프트 정의
    system_prompt = (
        f"""
        당신은 발표에 능숙하고 숙련된 전문가 AI입니다.
        발표자 정보: {user_profile.get('presenter', '미상')}, 발표 대상: {user_profile.get('audience', '미상')},
        발표 목적: {user_profile.get('purpose', '미상')}, 말투: {user_profile.get('tone', '미상')},
        발표 시간: {user_profile.get('time', '미상')}분.
        각 슬라이드 내용을 기반으로 자연스럽고 발표 톤으로 2~3문장 대본을 작성해 주세요.
        이전 슬라이드의 흐름을 고려해서 자연스럽게 이어지도록 해주세요.
        핵심 내용을 강조하여 전달력이 명확해야 합니다.
        분량은 각 슬라이드당 한 문단으로 만들어야 합니다.
        """
    )
    
    # ✅ 슬라이드별 대본 생성
    generated_scripts_per_slide = _presentation_scripts(slides_dict, system_prompt)

    # ✅ 전체 발표 대본 (초안)
    raw_script_combined = "\n\n".join(generated_scripts_per_slide[slide_num] for slide_num in sorted(generated_scripts_per_slide))

    # ✅ 최종 대본 다듬기
    polished_script = _polish_final_script(user_profile, raw_script_combined)
    
    # 세션에 저장 및 디버깅 로그
    session['generated_script_text'] = polished_script
    print(f"Saved to session: generated_script_text ({len(polished_script)} chars)")
    print(f"Session keys: {list(session.keys())}")

    return render_template('script_result.html', 
                          generated_script=polished_script,
                          selected_keywords=[item for sublist in all_selected_keywords_by_slide for item in sublist],
                          all_selected_keywords_by_slide=all_selected_keywords_by_slide,
                          script_options=script_options,
                          user_profile=user_profile)
    
@app.route('/script_generator/record', methods=['GET', 'POST'])
def record_script():
    generated_script = session.get('generated_script_text')
    if not generated_script:
        flash('먼저 대본을 생성해주세요.', 'danger')
        return redirect(url_for('record_script', reset=True))

    if request.method == 'POST':
        if 'audio_data' not in request.files:
            flash('오디오 파일이 업로드되지 않았습니다.', 'danger')
            return redirect(url_for('record_script'))

        audio_file = request.files['audio_data']
        
        if audio_file.filename == '':
            flash('오디오 파일을 선택해주세요.', 'danger')
            return redirect(url_for('record_script'))

        if not allowed_audio_file(audio_file.filename):
            flash('지원하지 않는 오디오 형식입니다. .mp3, .wav, .webm, .m4a 파일을 업로드해주세요.', 'danger')
            return redirect(url_for('record_script'))

        filepath = None
        try:
            original_filename = secure_filename(audio_file.filename)
            file_extension = os.path.splitext(original_filename)[1]
            unique_filename = str(uuid.uuid4()) + file_extension
            filepath = os.path.join(app.config['UPLOAD_FOLDER'], unique_filename)
            
            audio_file.save(filepath)
            
            # 파일 크기 검증
            file_size = os.path.getsize(filepath)
            if file_size > app.config['MAX_AUDIO_FILE_SIZE']:
                os.remove(filepath)
                flash(f'오디오 파일 크기가 너무 큽니다. 최대 {app.config["MAX_AUDIO_FILE_SIZE"] / (1024*1024):.0f}MB까지 업로드 가능합니다.', 'danger')
                return redirect(url_for('record_script'))
            
            if not os.path.exists(filepath):
                flash('오디오 파일 저장 중 오류가 발생했습니다. 다시 시도해주세요.', 'danger')
                return redirect(url_for('record_script'))

            session['recorded_audio_filename'] = unique_filename
            session['recorded_audio_original_filename'] = original_filename
            
            flash('녹음 파일이 성공적으로 업로드되었습니다. 이제 AI 평가를 시작할 수 있습니다.', 'success')
            return redirect(url_for('evaluation_result'))
            
        except Exception as e:
            flash(f'오디오 파일 업로드 중 오류가 발생했습니다: {e}', 'danger')
            print(f"Error during audio file saving: {e}")
            if filepath and os.path.exists(filepath):
                os.remove(filepath)
            return redirect(url_for('record_script'))

    return render_template('record_script.html', generated_script=generated_script)

# --- AI 평가 결과 페이지 (Whisper ASR 및 유사도 계산 통합) ---
@app.route('/script_generator/evaluation')
def evaluation_result():
    recorded_audio_filename = session.get('recorded_audio_filename')
    generated_script = session.get('generated_script_text')
    user_profile = session.get('user_profile')

    if not recorded_audio_filename or not generated_script:
        flash('평가에 필요한 정보(녹음 파일 또는 대본)가 부족합니다. 다시 시도해주세요.', 'danger')
        return redirect(url_for('user_info_form', reset=True))

    audio_filepath = os.path.join(app.config['UPLOAD_FOLDER'], recorded_audio_filename)
    wav_filepath = os.path.join(app.config['UPLOAD_FOLDER'], f"{os.path.splitext(recorded_audio_filename)[0]}.wav")

    recognized_text = "음성 인식 실패"
    fluency_score = 0.0

    try:
        # ✅ 오디오 파일을 WAV로 변환
        if not convert_to_wav(audio_filepath, wav_filepath):
            raise Exception("오디오 파일 변환 실패")

        # ✅ Whisper를 사용하여 음성 인식
        if pipe: # Whisper 모델이 성공적으로 로드되었는지 확인
            result = pipe(wav_filepath)
            recognized_text = result['text'].strip()
            # 텍스트가 너무 짧거나 인식되지 않았을 경우 처리
            if len(recognized_text) < 10: # 임의의 최소 길이
                recognized_text = "음성 인식이 제대로 되지 않았습니다. 명확하게 다시 녹음해보세요."
        else:
            raise Exception("Whisper 모델이 로드되지 않았습니다.")

        # ✅ 대본과 인식된 텍스트 간의 유사도 계산
        # 대본의 불필요한 마크다운 제거 및 공백 정규화
        cleaned_generated_script = generated_script.replace('**', '').replace('---', '').replace(':', '').strip()
        
        fluency_score = calc_similarity(recognized_text, cleaned_generated_script)

    except Exception as e:
        flash(f'음성 처리 및 평가 중 오류가 발생했습니다: {e}', 'danger')
        print(f"Error during ASR or similarity calculation: {e}")
        # 오류 발생 시 임시 파일 정리
        if os.path.exists(wav_filepath):
            os.remove(wav_filepath)
        # return redirect(url_for('record_script')) # 오류 시 다시 녹음 페이지로 이동
        # 오류가 발생하더라도 평가 페이지는 보여주되, 점수를 낮게 설정
        recognized_text = f"음성 처리 중 오류 발생: {e}"
        fluency_score = 0.0
    finally:
        # ✅ 사용이 끝난 WAV 파일 삭제
        if os.path.exists(wav_filepath):
            os.remove(wav_filepath)
        # 원본 오디오 파일도 평가 완료 후 삭제 (선택 사항)
        if os.path.exists(audio_filepath):
            os.remove(audio_filepath)


    # --- AI 평가 더미 로직 (나머지 점수) ---
    evaluation_scores = {
        'pronunciation': random.randint(70, 95),
        'speed': random.randint(60, 90),
        'emphasis': random.randint(50, 85),
        'confidence': random.randint(75, 98),
        'fluency': fluency_score # ✅ 실제 유사도 점수 반영
    }

    overall_score = sum(evaluation_scores.values()) / len(evaluation_scores)

    feedback_messages = []
    if evaluation_scores['pronunciation'] < 80:
        feedback_messages.append("일부 단어의 발음이 불분명할 수 있습니다. 반복 연습이 필요합니다.")
    if evaluation_scores['speed'] < 70:
        feedback_messages.append("발표 속도가 다소 빠르거나 느릴 수 있습니다. 청중의 이해를 위해 속도 조절을 고려해보세요.")
    elif evaluation_scores['speed'] > 85:
        feedback_messages.append("발표 속도가 적절합니다. 다만, 중요한 부분에서는 약간의 속도 변화를 주면 더욱 좋습니다.")
    if evaluation_scores['emphasis'] < 70:
        feedback_messages.append("핵심 메시지 강조가 부족할 수 있습니다. 중요한 단어에 힘을 주어 말하는 연습을 해보세요.")
    if evaluation_scores['confidence'] < 80:
        feedback_messages.append("억양에서 자신감이 부족하게 느껴질 수 있습니다. 확신을 가지고 말하는 연습을 해보세요.")
    
    # ✅ 실제 유사도 기반 피드백
    if fluency_score < 60:
        feedback_messages.append(f"대본과의 일치율이 낮습니다 ({fluency_score}%). 대본을 더 정확하게 읽는 연습이 필요합니다.")
    elif fluency_score < 80:
        feedback_messages.append(f"대본과의 일치율이 보통입니다 ({fluency_score}%). 몇몇 부분에서 대본과 다르게 발화된 부분이 있습니다.")
    else:
        feedback_messages.append(f"대본과의 일치율이 매우 높습니다 ({fluency_score}%). 발화가 대본과 잘 일치합니다!")
    
    if not feedback_messages:
        feedback_messages.append("전반적으로 훌륭한 발표였습니다! 계속해서 연습하시면 더욱 완벽해질 것입니다.")

    return render_template('evaluation_result.html',
                           overall_score=round(overall_score, 1),
                           evaluation_scores=evaluation_scores,
                           feedback_messages=feedback_messages,
                           generated_script=generated_script,
                           user_profile=user_profile,
                           recognized_text=recognized_text # 인식된 텍스트도 템플릿에 전달하여 디버깅/확인용으로 표시 가능
                           )


# --- 데이터베이스 생성 (최초 1회만 실행) ---
if __name__ == '__main__':
    with app.app_context():
       db.create_all()
    app.run(debug=True)
